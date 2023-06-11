import time
from sanic import Sanic
from sanic.response import text, file
from sanic.log import logger
from sanic import response
from sanic_cors import CORS

from pdfminer.high_level import extract_text
from pptx import Presentation

import os
import uuid
import bisect
import json
import re
from nltk.tokenize import word_tokenize
import nltk
from nltk.corpus import stopwords

from io import BytesIO

import spacy
import transformers


app = Sanic("Summarizer")
CORS(app)

model_name = "facebook/bart-base"

tokenizer = transformers.AutoTokenizer.from_pretrained(model_name)
model = transformers.AutoModelForSeq2SeqLM.from_pretrained(model_name)

nlp = spacy.load("en_core_web_sm")
nltk.download('stopwords')

def read_pdf(file_path):
    text = extract_text(file_path)
    return text

def read_pptx(file_path):
  prs = Presentation(file_path)
  text = ""
  for slide in prs.slides:
    for shape in slide.shapes:
      if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
          for run in paragraphs.runs:
            text += run.text

  return text

def clean_text(text):
  stop_words = set(stopwords.words('english'))

  text = text.lower()

  word_tokens = word_tokenize(text)

  filtered_text = [word for word in word_tokens if word.isalnum() and not word in stop_words]
  
  text = ' '.join(filtered_text)

  return text.strip()

def summarize_pdf(text): 
    doc = nlp(text)
    sentences = [clean_text(sent.text) for sent in doc.sents]
    processed_text = "\n".join(sentences)

    max_chunk_size = 2048
    text_chunks = [processed_text[i: i + max_chunk_size] for i in range(0, len(processed_text), max_chunk_size)]
    summaries = []

    for chunk in text_chunks:
         inputs = tokenizer.encode(chunk, return_tensors="pt", truncation=True)
         summary_ids = model.generate(inputs, num_beams=4, max_length=100, min_length=10, length_penalty=2.0, early_stopping=True)
         summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
         summaries.append(summary)

    completed_summary = " ".join(summaries)

    return completed_summary

# Determine whether to use 1024, 2048, or
# 4096 max_length to optimize performance
def determine_max_length(text):
    text_range = [1024, 2048, 4096]
    found_length = bisect.bisect(text_range, len(text))
    return text_range[found_length - 1]


def remove_non_latin1_chars(text):
    return text.encode('latin-1', 'ignore').decode('latin-1')


def write_summarization(file_path, text, file):
    if not os.path.exists('output'):
        os.makedirs('output')
    output_path = os.path.join("output", file_path)

    f = open(output_path, "w")
    f.write(json.dumps({ 'text': text, 'name': file.name.replace(".pdf", "").replace(".pptx", ""), 'date': round(time.time() * 1000) }))
    f.close()


def pdf_in_request(request):
    if not request.files or 'pdf_file' not in request.files:
        return 'No pdf file uploaded!'

    pdf_file = request.files.get('pdf_file')
    file_type = pdf_file.name.rsplit(".", 1)[1].lower()

    if not pdf_file or file_type != "pdf":
        return 'File uploaded is not a pdf file'

    return pdf_file

def pptx_in_request(request):
  if not request.files or 'pptx_file' not in request.files:
    return 'No pptx uploaded'

  pptx_file = request.files.get('pptx_file')
  file_type = pptx_file.name.rsplit(".", 1)[1].lower()

  if not pptx_file or file_type != "pptx":
   return "File uploaded is not a pptx file.."

  return pptx_file



@app.route('/summarize', methods=['POST'])
async def summarize(request):
    is_pdf_in_request = pdf_in_request(request)
    if type(is_pdf_in_request) == str:
        return response.json({"message": is_pdf_in_request}, 400)

    pdf_file = is_pdf_in_request

    try:
        text = read_pdf(BytesIO(pdf_file.body))
        summary = summarize_pdf(text)

        unique_id = str(uuid.uuid4())
        output_text_path = f"{unique_id}.json"

        write_summarization(output_text_path, summary, pdf_file)

        return response.json({"id": unique_id}, 200)

    except Exception as e:
        logger.error(f"Error during processing: {e}")
        return response.json({"message": "An error occured durring processing"}, 500)


@app.route('/summarization/<id>', methods=['GET'])
async def get_pdf(request, id):
    try:
        path = f"output/{id}.json"
        f = open(path, "r")
        info = json.loads(f.read())
        f.close()
        return response.json({"text": info["text"], "name": info["name"], "date": info["date"]})

    except FileNotFoundError as e:
        return response.json({"message": "File not found. It may have been deleted"}, 404)


@app.route('/pdf-length', methods=['GET'])
async def pdf_length(request):
    is_pdf_in_request = pdf_in_request(request)
    if type(is_pdf_in_request) == str:
        return response.json({"message": is_pdf_in_request}, 400)

    length = len(read_pdf(is_pdf_in_request))

    return response.json({"length": length})

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=3000)
